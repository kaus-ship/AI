from pptx import Presentation

def extract_ppt_text(path):

    prs = Presentation(path)

    slides = []

    for slide_no, slide in enumerate(prs.slides):

        text = ""

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

        slides.append(
            {
                "source": "PPT",
                "reference": f"Slide {slide_no+1}",
                "text": text
            }
        )

    return slides